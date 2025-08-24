import lxml.etree # type: ignore
import pptx, pptx.opc.constants, pptx.media, pptx.shapes.picture, pptx.slide
from model.entity.screen import Screen, MainScreen, MazeScreen, DialogScreen
from model.entity.state import State
from model.enum.asset.music import Music
from model.enum.asset.sound import Sound
from model.enum.button import Button
from model.enum.cell import Cell
from model.enum.effect import Effect
from util.logger import Logger


class Composer:

    def __init__(self, assets: dict[str, str], rows: float, columns: float) -> None:
        self.presentation = None
        self.width = None
        self.height = None
        self.size = None
        self.rows = rows
        self.columns = columns
        self.assets = assets
        for key, value in assets.items():
            if key.startswith('Image.'):
                assert value.endswith('.png'), f'Invalid image asset {key}={value}'
            elif key.startswith('Music.'):
                assert value.endswith('.mp3'), f'Invalid music asset {key}={value}'
            elif key.startswith('Sound.'):
                assert value.endswith('.wav'), f'Invalid sound asset {key}={value}'

    def prepare(self, count: int) -> None:
        self.presentation = pptx.Presentation()
        self.width = self.presentation.slide_width or 0
        self.height = self.presentation.slide_height or 0
        self.size = min(
            self.width / self.rows,
            self.height / self.columns
        )
        self.__prepare()
        for index in range(count):
            Logger.rewrite(f'Working: #{index} slides allocated')
            self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        Logger.write(f'Done: #{index} slides allocated')

    def __prepare(self) -> None:
        showpr_xml = '''
            <root
                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
            >
                <p:showPr loop="1" showNarration="1">
                    <p:kiosk/>
                    <p:sldAll/>
                    <p:penClr>
                        <a:prstClr val="red"/>
                    </p:penClr>
                    <p:extLst>
                        <p:ext uri="{EC167BDD-8182-4AB7-AECC-EB403E3ABB37}">
                            <p14:laserClr>
                                <a:srgbClr val="FF0000"/>
                            </p14:laserClr>
                        </p:ext>
                        <p:ext uri="{2FDB2607-1784-4EEB-B798-7EB5836EED8A}">
                        <p14:showMediaCtrls val="1"/>
                        </p:ext>
                    </p:extLst>
                </p:showPr>
            </root>
        '''
        showpr_element = lxml.etree.fromstring(showpr_xml)
        showpr_element = showpr_element.xpath('.//p:showPr', namespaces=showpr_element.nsmap)[0]
        presProps = self.presentation.part.part_related_by(pptx.opc.constants.RELATIONSHIP_TYPE.PRES_PROPS)
        element = lxml.etree.fromstring(presProps.blob)
        element.insert(0, showpr_element)
        presProps.blob = lxml.etree.tostring(element)

    def compose(self, screens: dict[str, tuple[Screen, dict[Button, str]]], limit: int = 1000000) -> pptx.presentation.Presentation:
        hashes = list(screens.keys())
        for index, context in enumerate(screens.values()):
            if index >= limit:
                break
            Logger.rewrite(f'Working: #{index} slides processed')
            screen, screen_actions = context
            slide = self.presentation.slides[index]
            slide_actions = self.__compose(slide, screen)
            if len(screen_actions) != len(slide_actions):
                raise ValueError(f'Action slides {len(screen_actions)} and buttons {len(slide_actions)} count mismatch')
            for action in screen_actions.keys():
                index = hashes.index(screen_actions[action])
                target = self.presentation.slides[index]
                button = slide_actions[action]
                button.click_action.target_slide = target
        Logger.write(f'Done: #{index} slides processed')
        return self.presentation

    def __compose(self, slide: pptx.slide.Slide, screen: Screen) -> dict[Button, pptx.shapes.picture.Picture]:
        if screen.music != Music.NONE:
            audio = slide.shapes.add_movie(
                self.assets[f'Music[{screen.music}]'],
                0,
                0,
                self.size,
                self.size,
                None,
                mime_type='audio/mp3'
            )
            self.set_autoplay(audio)
        if screen.sound != Sound.NONE:
            self.add_transition(slide, screen.sound, screen.effect)
        else:
            self.add_transition(slide, Sound.NONE, screen.effect)
        slide.shapes.add_picture(
            self.assets['Image.Menu[BACKGROUND]'],
            0,
            0,
            self.width,
            self.height
        )
        methods = {
            MainScreen.__name__: self.__compose_main,
            MazeScreen.__name__: self.__compose_maze,
            DialogScreen.__name__: self.__compose_dialog,
        }
        screen_type = screen.__class__.__name__
        if screen_type in methods:
            return methods[screen_type](slide, screen)
        else:
            raise ValueError(f'Unhandled screen {screen_type}')

    def __compose_main(self, slide: pptx.slide.Slide, screen: MainScreen) -> dict[Button, pptx.shapes.picture.Picture]:
        for corner in range(4):
            for tile in range(3):
                for side in range(2):
                    tile_offset = [
                        self.size * (2 + (tile if side == 0 else 0)),
                        self.size * (2 + (tile if side == 1 else 0))
                    ]
                    if corner % 2 == 1:
                        tile_offset[0] = self.width - tile_offset[0] - self.size
                    if corner // 2 == 1:
                        tile_offset[1] = self.height - tile_offset[1] - self.size
                    picture = slide.shapes.add_picture(
                        self.assets[f'Image.Maze[{Cell.WALL}]'],
                        tile_offset[0],
                        tile_offset[1],
                        self.size,
                        self.size
                    )
                    if tile == 0:
                        break
                    else:
                        self.set_transparency(picture, 1 - tile / 3)
        logo_size = (self.size * 10, self.size * 4)
        logo_offset = (
            (self.width - logo_size[0]) / 2,
            (self.height - logo_size[1]) / 2,
        )
        slide.shapes.add_picture(
            self.assets['Image.Menu[LOGO]'],
            logo_offset[0],
            logo_offset[1],
            logo_size[0],
            logo_size[1]
        )
        credits_size = (self.size * 4, self.size / 2)
        credits_offset = (
            self.width - credits_size[0],
            self.height - credits_size[1],
        )
        slide.shapes.add_picture(
            self.assets['Image.Menu[CREDITS]'],
            credits_offset[0],
            credits_offset[1],
            credits_size[0],
            credits_size[1]
        )
        return {
            Button.PLAY: slide.shapes.add_picture(
                self.assets[f'Image.Button[{Button.PLAY}]'],
                (self.width - self.size * 3) / 2,
                logo_offset[1] + logo_size[1] + self.size * 2,
                self.size * 3,
                self.size
            ),
        }

    def __compose_maze(self, slide: pptx.slide.Slide, screen: MazeScreen) -> dict[Button, pptx.shapes.picture.Picture]:
        return self.__paint_state(slide, screen.state)

    def __compose_dialog(self, slide: pptx.slide.Slide, screen: DialogScreen) -> dict[Button, pptx.shapes.picture.Picture]:
        self.__paint_state(slide, screen.state)
        picture = slide.shapes.add_picture(
            self.assets['Image.Menu[BACKGROUND]'],
            0,
            0,
            self.width,
            self.height
        )
        self.set_transparency(picture, 0.5)
        dialog_size = (self.size * 16, self.size * 4)
        dialog_offset = (
            (self.width - dialog_size[0]) / 2,
            (self.height - dialog_size[1]) * 2 / 3,
        )
        slide.shapes.add_picture(
            self.assets['Image.Menu[DIALOG]'],
            dialog_offset[0],
            dialog_offset[1],
            dialog_size[0],
            dialog_size[1]
        )
        textbox = slide.shapes.add_textbox(
            dialog_offset[0] + self.size,
            dialog_offset[1] + self.size,
            dialog_size[0] - self.size * 2,
            dialog_size[1] - self.size * 2
        ).text_frame
        textbox.text = self.assets[f'Message[{screen.message}]']
        textbox.word_wrap = True
        textbox.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
        textbox.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
        textbox.paragraphs[0].font.name = 'Courier New'
        textbox.paragraphs[0].font.size = pptx.util.Pt(24)
        textbox.paragraphs[0].font.color.rgb = pptx.dml.color.RGBColor(255, 255, 255)
        return {
            Button.CLOSE: slide.shapes.add_picture(
                self.assets[f'Image.Button[{Button.CLOSE}]'],
                dialog_offset[0] + dialog_size[0] - self.size,
                dialog_offset[1],
                self.size,
                self.size
            ),
        }

    def __paint_state(self, slide: pptx.slide.Slide, state: State) -> dict[Button, pptx.shapes.picture.Picture]:
        offset = (
            self.size * 5,
            self.size * 1
        )
        for top, line in enumerate(state.level.grid):
            for left, cell in enumerate(line):
                if cell == Cell.SPACE:
                    continue
                distance = max(
                    abs(state.position.left - left),
                    abs(state.position.top - top)
                )
                transparency = max(0.0, min(1.0, 1 - (distance - 1) / 3))
                if transparency <= 0:
                    continue
                picture = slide.shapes.add_picture(
                    self.assets[f'Image.Maze[{cell}]'],
                    left * self.size + offset[0],
                    top * self.size + offset[1],
                    self.size,
                    self.size
                )
                self.set_transparency(picture, transparency)
        slide.shapes.add_picture(
            self.assets[f'Image.Maze[{Cell.WIZARD}]'],
            state.position.left * self.size + offset[0],
            state.position.top * self.size + offset[1],
            self.size,
            self.size
        )
        for gems in range(3):
            if state.inventory.gems <= gems:
                image = self.assets[f'Image.Item[{Cell.GEM}].Missing']
            else:
                image = self.assets[f'Image.Item[{Cell.GEM}].Having']
            slide.shapes.add_picture(
                image,
                self.size * (1 + gems),
                self.size,
                self.size,
                self.size
            )
        for keys in range(1):
            if state.inventory.keys <= keys:
                continue
            slide.shapes.add_picture(
                self.assets[f'Image.Item[{Cell.KEY}].Having'],
                self.size * (1 + keys),
                self.size * 3,
                self.size,
                self.size
            )
        return {
            Button.UP: slide.shapes.add_picture(
                self.assets[f'Image.Button[{Button.UP}]'],
                self.size * 2,
                self.height - self.size * 4,
                self.size,
                self.size
            ),
            Button.LEFT: slide.shapes.add_picture(
                self.assets[f'Image.Button[{Button.LEFT}]'],
                self.size * 1,
                self.height - self.size * 3,
                self.size,
                self.size
            ),
            Button.RIGHT: slide.shapes.add_picture(
                self.assets[f'Image.Button[{Button.RIGHT}]'],
                self.size * 3,
                self.height - self.size * 3,
                self.size,
                self.size
            ),
            Button.DOWN: slide.shapes.add_picture(
                self.assets[f'Image.Button[{Button.DOWN}]'],
                self.size * 2,
                self.height - self.size * 2,
                self.size,
                self.size
            ),
        }

    def set_transparency(self, picture: pptx.shapes.picture.Picture, transparency: float) -> None:
        blip = picture.element.blipFill.blip
        blip.insert(
            0,
            lxml.etree.Element(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}alphaModFix',
                {'amt': str(int(transparency * 100000))}
            )
        )

    def set_autoplay(self, media: pptx.shapes.picture.Movie) -> None:
        media_id = media.element.xpath('.//p:cNvPr')[0].attrib['id']
        anchor_node = media.element.getparent().getparent().getparent().xpath(f'.//p:timing//p:video//p:spTgt[@spid="{media_id}"]')[0]
        media_node = anchor_node.getparent().getparent()
        media_node.set('numSld', '999999')
        element = media_node.xpath('.//p:cTn', namespaces=media_node.nsmap)[0]
        element.set('repeatCount', 'indefinite')
        element = media_node.xpath('.//p:cond', namespaces=media_node.nsmap)[0]
        element.set('delay', '0')

    def add_transition(self, slide: pptx.slide.Slide, sound: Sound, effect: Effect) -> None:
        sound_xml = ''
        if sound != Sound.NONE:
            audio = pptx.media.Video.from_path_or_file_like(self.assets[f'Sound[{sound}]'], 'audio/wav')
            media_part = slide.part.package.get_or_add_media_part(audio)
            slide.part.relate_to(media_part, pptx.opc.constants.RELATIONSHIP_TYPE.MEDIA)
            audio_id = slide.part.relate_to(media_part, pptx.opc.constants.RELATIONSHIP_TYPE.AUDIO)
            sound_xml = f'''
                <p:sndAc>
                    <p:stSnd>
                        <p:snd r:embed="{audio_id}" name="media{audio_id}.wav"/>
                    </p:stSnd>
                </p:sndAc>
            '''
        effect_xml = ''
        if effect != Effect.NONE:
            if effect == Effect.FADE:
                effect_xml = '<p:fade/>'
            elif effect == Effect.WIPE_UP:
                effect_xml = '<p:wipe dir="u"/>'
            else:
                raise ValueError(f'Unhandled effect {effect}')
        content_xml = f'''
            <root
                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"
                xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
            >
                <mc:AlternateContent>
                    <mc:Choice Requires="p14">
                        <p:transition advClick="0" p14:dur="100">
                            {effect_xml}
                            {sound_xml}
                        </p:transition>
                    </mc:Choice>
                    <mc:Fallback>
                        <p:transition advClick="0">
                            {effect_xml}
                            {sound_xml}
                        </p:transition>
                    </mc:Fallback>
                </mc:AlternateContent>
            </root>
        '''
        content_element = lxml.etree.fromstring(content_xml)
        content_element = content_element.xpath('.//mc:AlternateContent', namespaces=content_element.nsmap)[0]
        slide.element.insert(-1, content_element)
